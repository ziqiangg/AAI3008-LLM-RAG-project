# app/backend/services/ingestion.py
import os
import re
from typing import Callable, List, Optional, Dict, Any
from collections import defaultdict

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

from langchain_core.documents import Document as LCDocument
from langchain_community.document_loaders import (
    PDFPlumberLoader,
    Docx2txtLoader,
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from sqlalchemy.orm import Session as DBSession
from pptx import Presentation

from app.backend.config import Config
from app.backend.models import Document, DocumentChunk
from app.backend.services import classification

# ── Singleton: model loaded once at startup ───────────────────────────────
_embeddings: Optional[HuggingFaceEmbeddings] = None


def get_embeddings() -> HuggingFaceEmbeddings:
    """Lazy-load the HuggingFace embedding model once."""
    global _embeddings
    if _embeddings is None:
        _embeddings = HuggingFaceEmbeddings(
            model_name=Config.EMBEDDING_MODEL,
            model_kwargs={"device": "cpu"},
            encode_kwargs={"normalize_embeddings": True},
        )
    return _embeddings


# ──────────────────────────────────────────────────────────────────────────
# Semantic chunking helpers
# ──────────────────────────────────────────────────────────────────────────

# Updated to support multilingual sentence splitting (English + Chinese punctuation)
_SENTENCE_SPLIT_RE = re.compile(r"(?<=[.!?。！？])\s*")
_MULTI_NL_RE = re.compile(r"\n{2,}")


def _looks_like_bullets(text: str) -> bool:
    """Detect bullet-heavy text (common in slides)."""
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    if len(lines) < 5:
        return False
    bulletish = 0
    for ln in lines:
        if ln.startswith(("•", "·", "-", "–", "—", "*")):
            bulletish += 1
        # numbered bullets
        if re.match(r"^\d+[\).\]]\s+", ln):
            bulletish += 1
    return bulletish / max(len(lines), 1) >= 0.35


def _split_into_units(text: str, max_unit_chars: int = 600) -> List[str]:
    """
    Split text into "semantic-ish" units (bullets/paragraphs, then sentences).
    Units should be reasonably sized for efficient embedding.
    Increased from 400 to 600 chars for better performance.
    """
    if not text:
        return []

    text = text.strip()

    # If it's bullet-heavy, treat each non-empty line as a unit, then pack.
    if _looks_like_bullets(text):
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        units: List[str] = []
        buf: List[str] = []
        buf_len = 0
        for ln in lines:
            if buf_len + len(ln) + 1 > max_unit_chars and buf:
                units.append(" ".join(buf).strip())
                buf, buf_len = [], 0
            buf.append(ln)
            buf_len += len(ln) + 1
        if buf:
            units.append(" ".join(buf).strip())
        return units

    # Otherwise: paragraph first
    paras = [p.strip() for p in _MULTI_NL_RE.split(text) if p.strip()]
    units: List[str] = []

    for p in paras:
        if len(p) <= max_unit_chars:
            units.append(p)
            continue

        # Split into sentences and pack
        sents = [s.strip() for s in _SENTENCE_SPLIT_RE.split(p) if s.strip()]
        buf: List[str] = []
        buf_len = 0
        for s in sents:
            if buf_len + len(s) + 1 > max_unit_chars and buf:
                units.append(" ".join(buf).strip())
                buf, buf_len = [], 0
            buf.append(s)
            buf_len += len(s) + 1
        if buf:
            units.append(" ".join(buf).strip())

    return units


def _semantic_chunk_single_doc(
    doc: LCDocument,
    *,
    max_chunk_chars: int = 1400,
    min_chunk_chars: int = 250,
    similarity_threshold: float = 0.55,
    debug: bool = False,
) -> List[LCDocument]:
    """
    Semantic chunking:
    - Split doc into units (bullets/paragraph/sentence packs)
    - Embed each unit (normalized embeddings already)
    - Start a new chunk when adjacent similarity drops below threshold
    - Also respect max_chunk_chars to avoid overly large chunks
    """
    text = (doc.page_content or "").strip()
    if not text:
        return []

    units = _split_into_units(text)
    if len(units) <= 1:
        return [doc]

    vectors = get_embeddings().embed_documents(units)

    # Cosine similarity = dot product since vectors normalized
    def dot(a, b) -> float:
        return sum(x * y for x, y in zip(a, b))

    chunks: List[LCDocument] = []
    cur_units = [units[0]]
    cur_len = len(units[0])

    for i in range(1, len(units)):
        sim = dot(vectors[i - 1], vectors[i])

        # break if topic shift OR chunk too big
        should_break = (sim < similarity_threshold) or (cur_len + len(units[i]) + 1 > max_chunk_chars)

        if should_break and cur_len >= min_chunk_chars:
            if debug:
                print(f"[semantic split] sim={sim:.2f} break at unit {i-1}->{i} meta={doc.metadata}")
            chunks.append(
                LCDocument(
                    page_content="\n\n".join(cur_units).strip(),
                    metadata=dict(doc.metadata or {}),
                )
            )
            cur_units = [units[i]]
            cur_len = len(units[i])
        else:
            cur_units.append(units[i])
            cur_len += len(units[i]) + 1

    if cur_units:
        chunks.append(
            LCDocument(
                page_content="\n\n".join(cur_units).strip(),
                metadata=dict(doc.metadata or {}),
            )
        )

    return chunks


def semantic_chunk_documents(
    lc_docs: List[LCDocument],
    *,
    max_chunk_chars: int = 1400,
    min_chunk_chars: int = 250,
    similarity_threshold: float = 0.55,
    debug: bool = False,
) -> List[LCDocument]:
    """Apply semantic chunking across a list of docs, preserving metadata."""
    all_chunks: List[LCDocument] = []
    for d in lc_docs:
        all_chunks.extend(
            _semantic_chunk_single_doc(
                d,
                max_chunk_chars=max_chunk_chars,
                min_chunk_chars=min_chunk_chars,
                similarity_threshold=similarity_threshold,
                debug=debug,
            )
        )
    return all_chunks


# ──────────────────────────────────────────────────────────────────────────
# Loading
# ──────────────────────────────────────────────────────────────────────────

def _load_pdf_with_pymupdf(file_path: str) -> List[LCDocument]:
    """Load PDF using PyMuPDF (fitz) - better for mathematical content."""
    try:
        doc = fitz.open(file_path)
        lc_docs = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            # Extract text with layout preservation
            text = page.get_text("text", sort=True)
            
            # Clean up common PDF extraction artifacts
            text = text.replace('\x00', '')  # Remove null bytes
            text = re.sub(r'\s+', ' ', text)  # Normalize whitespace
            text = text.strip()
            
            if text:  # Only add non-empty pages
                lc_docs.append(
                    LCDocument(
                        page_content=text,
                        metadata={
                            "source": file_path,
                            "page": page_num + 1,
                            "total_pages": len(doc),
                        }
                    )
                )
        
        doc.close()
        return lc_docs if lc_docs else [LCDocument(page_content="", metadata={"source": file_path})]
    
    except Exception as e:
        raise ValueError(f"PyMuPDF loading failed: {e}")


def _load_text_with_encoding(file_path: str) -> List[LCDocument]:
    """Load text file with UTF-8 encoding and fallback handling."""
    encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding, errors='strict') as f:
                content = f.read()
            
            # Ensure content is valid UTF-8 for storage
            content.encode('utf-8')
            
            return [LCDocument(
                page_content=content,
                metadata={"source": file_path, "encoding": encoding}
            )]
        except (UnicodeDecodeError, UnicodeEncodeError):
            continue
    
    # Final fallback: read with errors='replace' to avoid complete failure
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        return [LCDocument(
            page_content=content,
            metadata={"source": file_path, "encoding": "utf-8-fallback"}
        )]
    except Exception as e:
        raise ValueError(f"Text file loading failed: {e}")


def _load_pptx_with_python_pptx(file_path: str) -> List[LCDocument]:
    """Load PPTX using python-pptx library directly (no API calls)."""
    try:
        prs = Presentation(file_path)
        docs = []
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # Ensure text is properly decoded
                    text = shape.text.strip()
                    # Normalize text to ensure valid UTF-8
                    try:
                        text = text.encode('utf-8', errors='ignore').decode('utf-8')
                    except:
                        pass
                    slide_text_parts.append(text)
            
            # Combine all text for this slide
            slide_text = "\n".join(slide_text_parts).strip()
            
            if slide_text:  # Only add non-empty slides
                docs.append(
                    LCDocument(
                        page_content=slide_text,
                        metadata={
                            "source": file_path,
                            "slide_number": slide_idx,
                            "total_slides": len(prs.slides),
                        }
                    )
                )
        
        return docs if docs else [LCDocument(page_content="", metadata={"source": file_path})]
    except Exception as e:
        raise ValueError(f"PPTX loading failed: {e}")

def load_document(file_path: str) -> List[LCDocument]:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pptx":
        return _load_pptx_with_python_pptx(file_path)

    # Handle text files with explicit UTF-8 encoding and fallback
    if ext == ".txt":
        return _load_text_with_encoding(file_path)
    
    # Handle PDFs with PyMuPDF (better for mathematical content)
    if ext == ".pdf":
        if PYMUPDF_AVAILABLE:
            try:
                return _load_pdf_with_pymupdf(file_path)
            except Exception as e:
                print(f"PyMuPDF failed ({e}), falling back to PDFPlumber...")
                return PDFPlumberLoader(file_path).load()
        else:
            return PDFPlumberLoader(file_path).load()

    loader_map = {
        ".docx": Docx2txtLoader,
    }

    loader_cls = loader_map.get(ext)
    if not loader_cls:
        raise ValueError(
            f"Unsupported file extension '{ext}'. Allowed: {Config.ALLOWED_EXTENSIONS}"
        )

    return loader_cls(file_path).load()


# ──────────────────────────────────────────────────────────────────────────
# Fallback chunking (recursive)
# ──────────────────────────────────────────────────────────────────────────

def _split_with_recursive(
    docs: List[LCDocument],
    *,
    chunk_size: int,
    chunk_overlap: int,
    separators=None
) -> List[LCDocument]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=separators or ["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_documents(docs)


# ──────────────────────────────────────────────────────────────────────────
# PPTX: merge loader elements into one doc per slide
# ──────────────────────────────────────────────────────────────────────────

def _pptx_docs_by_slide(lc_docs: List[LCDocument]) -> List[LCDocument]:
    """
    Merge multiple document elements into one doc per slide.
    Now mainly used if we switch back to UnstructuredPowerPointLoader.
    With python-pptx loader, docs are already slide-based.
    """
    # If already slide-based (from python-pptx), just return
    if all(d.metadata.get("slide_number") for d in lc_docs if d.metadata):
        return lc_docs
    
    slides: Dict[int, List[str]] = defaultdict(list)
    slide_meta: Dict[int, Dict[str, Any]] = {}

    for d in lc_docs:
        md = d.metadata or {}
        slide_no = (
            md.get("page_number")
            or md.get("slide_number")
            or md.get("page")
            or md.get("slide")
            or 1
        )
        slide_no = int(slide_no)

        txt = (d.page_content or "").strip()
        if txt:
            slides[slide_no].append(txt)

        # keep first seen metadata as base, then update with any missing keys
        if slide_no not in slide_meta:
            slide_meta[slide_no] = dict(md)
        else:
            for k, v in md.items():
                if k not in slide_meta[slide_no]:
                    slide_meta[slide_no][k] = v

    merged: List[LCDocument] = []
    for slide_no in sorted(slides.keys()):
        slide_text = "\n".join(slides[slide_no]).strip()
        if not slide_text:
            continue
        md = dict(slide_meta.get(slide_no, {}))
        md["slide_number"] = slide_no
        merged.append(LCDocument(page_content=slide_text, metadata=md))

    return merged


# ──────────────────────────────────────────────────────────────────────────
# PDF: merge slide-like pages into windows so semantic chunking has context
# ──────────────────────────────────────────────────────────────────────────

def _is_slide_like_pdf(lc_docs: List[LCDocument]) -> bool:
    """
    Heuristic: slide PDFs have short pages with few punctuation marks.
    """
    if not lc_docs:
        return False

    samples = lc_docs[: min(len(lc_docs), 12)]
    lengths = [len((d.page_content or "").strip()) for d in samples]
    avg_len = sum(lengths) / max(len(lengths), 1)

    punct = 0
    chars = 0
    for d in samples:
        t = (d.page_content or "")
        punct += sum(t.count(x) for x in [".", "!", "?", ";", ":"])
        chars += len(t)

    punct_rate = punct / max(chars, 1)

    # Typical: avg page small and punctuation sparse
    return (avg_len < 700) and (punct_rate < 0.01)


def _merge_pdf_pages(lc_docs: List[LCDocument], window_pages: int = 3) -> List[LCDocument]:
    """
    Merge every N pages into one doc, keeping start/end page metadata.
    """
    merged: List[LCDocument] = []
    buf: List[str] = []
    md_acc: Dict[str, Any] = {}
    start_page = None
    end_page = None

    def flush():
        nonlocal buf, md_acc, start_page, end_page
        if not buf:
            return
        md_out = dict(md_acc)
        if start_page is not None:
            md_out["start_page"] = start_page
        if end_page is not None:
            md_out["end_page"] = end_page
        merged.append(LCDocument(page_content="\n\n".join(buf).strip(), metadata=md_out))
        buf, md_acc, start_page, end_page = [], {}, None, None

    for i, d in enumerate(lc_docs):
        t = (d.page_content or "").strip()
        md = d.metadata or {}

        p = md.get("page")
        if start_page is None and p is not None:
            start_page = p
        if p is not None:
            end_page = p

        if not md_acc:
            md_acc = dict(md)
        else:
            # keep base, fill missing keys only
            for k, v in md.items():
                if k not in md_acc:
                    md_acc[k] = v

        if t:
            buf.append(t)

        if (i + 1) % window_pages == 0:
            flush()

    flush()
    return merged


# ──────────────────────────────────────────────────────────────────────────
# Cleaning
# ──────────────────────────────────────────────────────────────────────────

def clean_slide_text(text: str) -> str:
    if not text:
        return ""
    lines = [ln.rstrip() for ln in text.splitlines()]
    cleaned = []
    for ln in lines:
        low = ln.strip().lower()
        if low in {"sit internal"}:
            continue
        cleaned.append(ln)

    out = "\n".join([ln for ln in cleaned if ln.strip() != ""]).strip()
    return out


def _dedouble_word_if_encoded(word: str) -> str:
    """
    Fix OCR/PDF extraction artifacts like:
      'CClloouudd' -> 'Cloud'
    WITHOUT breaking normal words like:
      'access' (should remain 'access', not 'aces')
    Strategy:
      If the word is even-length and mostly made of repeated pairs (aa bb cc ...),
      collapse pairs.
    """
    if not word:
        return word
    if len(word) < 6 or len(word) % 2 != 0:
        return word

    # Check pair structure
    pairs = [(word[i], word[i + 1]) for i in range(0, len(word), 2)]
    same_pairs = sum(1 for a, b in pairs if a == b)
    ratio = same_pairs / max(len(pairs), 1)

    # Only collapse if *most* pairs are duplicated
    if ratio >= 0.8:
        return "".join(a for a, _ in pairs)

    return word


def clean_pdf_extraction_noise(text: str) -> str:
    if not text:
        return ""
    
    # 0) Ensure text is valid UTF-8 and normalize encoding issues
    try:
        # Remove common mojibake patterns
        text = text.encode('utf-8', errors='ignore').decode('utf-8')
        # Replace common problematic characters
        text = text.replace('\ufffd', '')  # Remove replacement character
        text = text.replace('\x00', '')     # Remove null bytes
    except:
        pass

    # 1) Fix "doubled-letter words" safely (token-wise)
    def fix_token(tok: str) -> str:
        # split punctuation from word-ish core
        m = re.match(r"^(\W*)([A-Za-z]{6,})(\W*)$", tok)
        if not m:
            return tok
        pre, core, post = m.group(1), m.group(2), m.group(3)
        core2 = _dedouble_word_if_encoded(core)
        return f"{pre}{core2}{post}"

    tokens = re.split(r"(\s+)", text)  # keep whitespace
    tokens = [fix_token(t) if not t.isspace() else t for t in tokens]
    text = "".join(tokens)

    # 2) Fix spaced letters lines like "I P S" or "n l o"
    fixed_lines = []
    for line in text.splitlines():
        stripped = line.strip()
        if re.fullmatch(r"(?:[A-Za-z]\s+){2,}[A-Za-z]", stripped):
            fixed_lines.append(stripped.replace(" ", ""))
        else:
            fixed_lines.append(line)
    text = "\n".join(fixed_lines)

    # 3) Merge broken line wraps (common in PDFs)
    lines = [ln.rstrip() for ln in text.splitlines()]
    merged = []
    i = 0
    while i < len(lines):
        cur = lines[i].strip()
        if not cur:
            i += 1
            continue
        if i + 1 < len(lines):
            nxt = lines[i + 1].strip()
            if cur and nxt and (not re.search(r"[.!?:;]$", cur)) and re.match(r"^[a-z]", nxt):
                merged.append(cur + " " + nxt)
                i += 2
                continue
        merged.append(cur)
        i += 1

    return "\n".join(merged).strip()


# ──────────────────────────────────────────────────────────────────────────
# Type-aware splitting
# ──────────────────────────────────────────────────────────────────────────

def split_documents_by_type(lc_docs: List[LCDocument], file_ext: str) -> List[LCDocument]:
    """
    Adaptive chunking: Use semantic chunking for smaller/medium docs,
    fall back to fast recursive chunking for very large documents.
    """
    ext = f".{file_ext.lower().lstrip('.')}"
    
    # Calculate total document size
    total_chars = sum(len(d.page_content or "") for d in lc_docs)
    
    # Performance threshold: if document is too large, use simpler chunking
    LARGE_DOC_THRESHOLD = 150000  # ~150KB of text
    use_fast_chunking = total_chars > LARGE_DOC_THRESHOLD
    
    if use_fast_chunking:
        print(f"     ℹ️  Large document ({total_chars} chars), using fast recursive chunking")
        return _split_with_recursive(
            lc_docs,
            chunk_size=1000,
            chunk_overlap=100,
        )
    
    # --- PPTX (already slide-based, usually small) ---
    if ext == ".pptx":
        # No need to merge slides, already loaded per-slide
        return semantic_chunk_documents(
            lc_docs,
            max_chunk_chars=1200,  # Reduced for faster embedding
            min_chunk_chars=200,
            similarity_threshold=0.58,  # Slightly higher = fewer chunks
            debug=False,
        )

    # --- PDF ---
    if ext == ".pdf":
        # If it's a slide-like PDF, merge pages into windows first
        if _is_slide_like_pdf(lc_docs):
            merged_docs = _merge_pdf_pages(lc_docs, window_pages=3)
            return semantic_chunk_documents(
                merged_docs,
                max_chunk_chars=1400,
                min_chunk_chars=250,
                similarity_threshold=0.58,
                debug=False,
            )

        # normal PDFs: semantic chunk per page
        return semantic_chunk_documents(
            lc_docs,
            max_chunk_chars=1200,
            min_chunk_chars=250,
            similarity_threshold=0.58,
            debug=False,
        )

    # --- DOCX ---
    if ext == ".docx":
        return semantic_chunk_documents(
            lc_docs,
            max_chunk_chars=1200,
            min_chunk_chars=250,
            similarity_threshold=0.58,
            debug=False,
        )

    # --- TXT/default ---
    return semantic_chunk_documents(
        lc_docs,
        max_chunk_chars=1000,
        min_chunk_chars=200,
        similarity_threshold=0.58,
        debug=False,
    )


# ──────────────────────────────────────────────────────────────────────────
# Embedding
# ──────────────────────────────────────────────────────────────────────────

def embed_texts(texts: List[str]) -> List[List[float]]:
    return get_embeddings().embed_documents(texts)


# ──────────────────────────────────────────────────────────────────────────
# Ingestion pipeline
# ──────────────────────────────────────────────────────────────────────────

def run_ingestion_pipeline(
    db_session: DBSession,
    file_path: str,
    user_id: Optional[int] = None,
    subject: Optional[str] = None,
    progress_callback: Optional[Callable[..., None]] = None,
) -> Document:
    filename = os.path.basename(file_path)
    file_ext = os.path.splitext(filename)[1].lstrip(".")

    def report(progress: Optional[int] = None, status_message: Optional[str] = None, **details) -> None:
        if progress_callback is None:
            return
        payload = {key: value for key, value in details.items() if value is not None}
        if progress is not None:
            payload["progress"] = progress
        if status_message:
            payload["status_message"] = status_message
        if payload:
            progress_callback(**payload)

    report(20, "Reading document")

    # ?????? 1. Load ??????
    print(f"[1] Loading  : {filename}")
    lc_docs = load_document(file_path)
    print(f"     ??? {len(lc_docs)} page(s)/element-doc(s) loaded")
    report(30, f"Loaded {len(lc_docs)} page(s)")

    # ?????? 1b. Classify Document Subjects ??????
    print("[1b] Classifying document subjects...")
    report(40, "Classifying subject")
    embeddings_model = get_embeddings()

    # Extract sample from first few documents/pages
    content_sample = ""
    for doc in lc_docs[:5]:  # First 5 pages/elements
        content_sample += doc.page_content + "\n"
    content_sample = content_sample[:3000]  # Limit to 3000 chars

    # Classify subjects
    classified_subjects = classification.classify_document_subjects(
        content_sample=content_sample,
        embeddings_model=embeddings_model
    )

    # Extract subject names for document record
    document_subjects = [s['name'] for s in classified_subjects]
    print(f"     ??? Classified as: {', '.join(document_subjects)}")
    print(f"     ??? Confidence: {classified_subjects[0]['confidence']:.2f}")
    report(48, "Subject classified", subjects=document_subjects)

    # ?????? 2. Chunk ??????
    print("[2] Chunking : method=semantic_embedding")
    report(58, "Chunking document", subjects=document_subjects)
    chunks = split_documents_by_type(lc_docs, file_ext)
    print(f"     ??? {len(chunks)} raw chunks")

<<<<<<< Updated upstream
    # ── 2b. Clean + Filter (BEFORE EMBEDDING/STORING) ──
    MIN_CHUNK_LEN = 140  # slightly higher: removes more junk
=======
    # ?????? 2b. Clean + Filter (BEFORE EMBEDDING/STORING) ??????
    MIN_CHUNK_LEN = 30  # lowered: avoids rejecting short but valid content
>>>>>>> Stashed changes
    cleaned_chunks: List[LCDocument] = []

    for c in chunks:
        text = c.page_content or ""
        text = clean_slide_text(text)
        text = clean_pdf_extraction_noise(text)

        text = text.strip()
        if len(text) < MIN_CHUNK_LEN:
            continue

        c.page_content = text
        cleaned_chunks.append(c)

    chunks = cleaned_chunks
    print(f"     ??? {len(chunks)} cleaned chunks kept")
    report(68, "Chunks prepared", subjects=document_subjects, chunk_count=len(chunks))

    if not chunks:
        raise ValueError(
            "No usable chunks after cleaning/filtering. "
            "Try lowering MIN_CHUNK_LEN or check the document extraction quality."
        )

    # ?????? 3. Embed ??????
    print(f"[3] Embedding: {len(chunks)} chunks via {Config.EMBEDDING_MODEL}")
    report(76, "Generating embeddings", subjects=document_subjects, chunk_count=len(chunks))
    texts = [c.page_content for c in chunks]
    vectors = embed_texts(texts)

    if not vectors or not vectors[0]:
        raise RuntimeError("Embedding failed: got empty vectors.")
    print(f"     ??? dim={len(vectors[0])}")

    # ?????? 3b. Classify Chunk Topics ??????
    print(f"[3b] Classifying topics for {len(chunks)} chunks...")
    report(84, "Classifying chunk topics", subjects=document_subjects, chunk_count=len(chunks))
    chunk_classifications = []

    for idx, (chunk, vector) in enumerate(zip(chunks, vectors)):
        # Classify topics for this chunk
        topics = classification.classify_chunk_topics(
            chunk_content=chunk.page_content,
            chunk_embedding=vector,
            document_subjects=document_subjects
        )
        chunk_classifications.append(topics)

    print(f"     ??? Topic classification complete")

    # ?????? 4. Create Document row ??????
    doc = Document(
        user_id=user_id,
        filename=filename,
        file_path=os.path.abspath(file_path),
        file_type=file_ext,
        title=filename,
        subject=document_subjects,  # Array of classified subjects
        chunk_count=len(chunks),
    )
    db_session.add(doc)
    db_session.flush()
    print(f"[4] Document row created ??? id={doc.id}")
    report(90, "Saving document", document_id=doc.id, subjects=document_subjects, chunk_count=len(chunks))

    # ?????? 5. Insert chunks ??????
    # Keep params in one place so metadata matches real behavior
    semantic_params = {
        "similarity_threshold": 0.55,
        "max_chunk_chars": 1400,
        "min_chunk_chars": 300 if file_ext.lower() in ["pdf", "docx"] else 250,
    }

    progress_step = max(1, len(chunks) // 4) if chunks else 1

    for idx, (chunk, vector) in enumerate(zip(chunks, vectors)):
        # Get topic classification for this chunk
        topics = chunk_classifications[idx]

        # Determine dominant subject and topic for quick access
        dominant_subject = topics[0]['name'] if topics else document_subjects[0]
        dominant_topic = ""
        if topics and topics[0].get('topics'):
            top_topic = topics[0]['topics'][0]
            dominant_topic = f"{top_topic['name']}/{top_topic.get('subtopic', '')}".rstrip('/')

        db_session.add(
            DocumentChunk(
                document_id=doc.id,
                chunk_order=idx,
                content=chunk.page_content,
                embedding=vector,
                chunk_metadata={
                    "chunk_index": idx,
                    "total_chunks": len(chunks),
                    "chunking_method": "semantic_embedding",
                    "semantic_params": semantic_params,
                    "source": chunk.metadata,
                    "content_len": len(chunk.page_content or ""),
                    # NEW: Subject and topic classification
                    "subjects": topics,
                    "dominant_subject": dominant_subject,
                    "dominant_topic": dominant_topic,
                },
            )
        )

        if (idx + 1) == len(chunks) or (idx + 1) % progress_step == 0:
            stored_ratio = (idx + 1) / len(chunks)
            progress = min(98, 90 + int(stored_ratio * 8))
            report(progress, f"Saving chunks ({idx + 1}/{len(chunks)})", document_id=doc.id, subjects=document_subjects, chunk_count=len(chunks))

    db_session.commit()
    print(f"[5] Stored   : {len(chunks)} chunks ??? document_id={doc.id}")
    print("???  Ingestion complete.")
    report(100, "Indexed and ready", document_id=doc.id, subjects=document_subjects, chunk_count=len(chunks))

    # Preview first few chunks
    for i, c in enumerate(chunks[:5]):
        print(f"\n--- CHUNK {i} ---")
        print(c.page_content[:1200])
        print("LEN:", len(c.page_content))

    return doc

def _split_long_text(text: str, chunk_size: int) -> list[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=0,
        separators=["\n\n", "\n", " ", ""]
    )
    return splitter.split_text(text or "")


def chunk_sections_to_db(
    db_session,
    document_id: int,
    sections: list[dict],
    source_metadata: dict,
) -> tuple[int, list[dict]]:
    """
    Hybrid link ingestion:
    - split HTML into heading sections upstream
    - each section becomes a chunk
    - if a section is too long, split inside the section (fallback)
    - add subject/topic metadata using existing classification service

    Metadata schema aligns with your normal docs:
      - subjects: List[{"name": str, "confidence": float}, ...]
      - dominant_subject: str
      - dominant_topic: str (e.g. "Topic/Subtopic")
    """
    from app.backend.services import classification  # uses classify_document_subjects + classify_chunk_topics
    from app.backend.models import DocumentChunk
    from langchain.text_splitter import RecursiveCharacterTextSplitter

    emb = get_embeddings()

    max_chars = getattr(Config, "LINK_CHUNK_MAX_CHARS", 2200)
    min_chars = getattr(Config, "LINK_CHUNK_MIN_CHARS", 350)
    long_split = getattr(Config, "LINK_LONG_SECTION_SPLIT_SIZE", 1200)

    # ─────────────────────────────────────────────────────────
    # 1) Build candidate chunk texts (hybrid: section + fallback split)
    # ─────────────────────────────────────────────────────────
    def split_long_text(text: str) -> list[str]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=long_split,
            chunk_overlap=0,
            separators=["\n\n", "\n", " ", ""]
        )
        return splitter.split_text(text or "")

    candidate_chunks: list[dict] = []  # [{"title":..., "text":...}, ...]

    for sec in sections:
        title = (sec.get("title") or "Section").strip()
        text = (sec.get("text") or "").strip()
        if not text:
            continue

        parts = [text] if len(text) <= max_chars else split_long_text(text)

        # Merge small fragments within a section to avoid micro-chunks
        merged_parts = []
        buf = ""
        for p in parts:
            p = (p or "").strip()
            if not p:
                continue
            if not buf:
                buf = p
            elif len(buf) + len(p) + 2 <= max_chars:
                buf = buf + "\n\n" + p
            else:
                merged_parts.append(buf)
                buf = p
        if buf:
            merged_parts.append(buf)

        for chunk_text in merged_parts:
            chunk_text = chunk_text.strip()
            if len(chunk_text) < min_chars:
                continue
            candidate_chunks.append({"title": title, "text": chunk_text})

    if not candidate_chunks:
        return 0

    # ─────────────────────────────────────────────────────────
    # 2) Document-level subject classification (once)
    #    Use a sample from the first few sections for speed.
    # ─────────────────────────────────────────────────────────
    sample_text = "\n\n".join(
        [f"{c['title']}\n{c['text']}" for c in candidate_chunks[:3]]
    )
    
    # Returns: [{"name": "AI", "confidence": 0.82}, ...]
    doc_subject_results = classification.classify_document_subjects(
        content_sample=sample_text[:3000],
        embeddings_model=emb,  # works because it provides embed_query()
        threshold=getattr(Config, "SUBJECT_SIMILARITY_THRESHOLD", None)
    )
    if not doc_subject_results:
        doc_subject_results = [{"name": "General", "confidence": 1.0}]

    document_subject_names = [s["name"] for s in doc_subject_results]
    dominant_subject = doc_subject_results[0]["name"]

    # ─────────────────────────────────────────────────────────
    # 3) Insert chunks with chunk-level topic classification
    # ─────────────────────────────────────────────────────────
    total_chunks = len(candidate_chunks)
    order = 1
    count = 0

    for c in candidate_chunks:
        title = c["title"]
        chunk_text = c["text"]

        vec = emb.embed_query(chunk_text)

        # Topic classification within known document subjects
        topic_results = classification.classify_chunk_topics(
            chunk_content=chunk_text,
            chunk_embedding=vec,
            document_subjects=document_subject_names,
            threshold=getattr(Config, "TOPIC_SIMILARITY_THRESHOLD", None)
        )

        # Choose a dominant topic string like "Topic/Subtopic"
        dominant_topic = ""
        if topic_results:
            # pick best subject group, then best topic within it
            best_group = sorted(topic_results, key=lambda x: x.get("confidence", 0.0), reverse=True)[0]
            top_topics = best_group.get("topics") or []
            if top_topics:
                t0 = top_topics[0]
                if t0.get("subtopic"):
                    dominant_topic = f"{t0.get('name')}/{t0.get('subtopic')}"
                else:
                    dominant_topic = f"{t0.get('name')}"

        md = {
            "chunk_index": order,
            "total_chunks": total_chunks,
            "chunking_method": "hybrid_heading_sections",
            "content_len": len(chunk_text),

            # link-specific
            "source_type": "link",
            "url": (source_metadata or {}).get("url"),
            "section_title": title,

            # subject/topic metadata (aligned with your normal docs + extract_subject_context)
            "subjects": doc_subject_results,          # List[{"name","confidence"}]
            "dominant_subject": dominant_subject,     # string
            "dominant_topic": dominant_topic,         # string "Topic/Subtopic"
            "topic_matches": topic_results,           # optional detail for debugging/UI
        }

        dc = DocumentChunk(
            document_id=document_id,
            chunk_order=order,
            content=chunk_text,
            embedding=vec,
            chunk_metadata=md
        )
        db_session.add(dc)
        order += 1
        count += 1

    return count, doc_subject_results